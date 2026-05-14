from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# --- Theme ---
SLIDE_W = Inches(11.69)
SLIDE_H = Inches(8.27)

COLORS = {
    "bg": RGBColor(252, 250, 255),
    "primary": RGBColor(140, 92, 255),      # #8C5CFF
    "secondary": RGBColor(198, 166, 255),   # #C6A6FF
    "deep_text": RGBColor(47, 35, 64),      # #2F2340
    "soft_text": RGBColor(123, 108, 142),   # #7B6C8E
    "card": RGBColor(248, 244, 252),        # #F8F4FC
    "pink": RGBColor(246, 223, 255),        # #F6DFFF
    "line": RGBColor(226, 214, 244),
    "white": RGBColor(255, 255, 255),
}

FONT_KR = "Pretendard"
FONT_FALLBACK = "Noto Sans KR"

IMAGE_FILES = {
    "ticket": "选票.png",
    "date": "观展日期选择.png",
    "order": "订单审核.png",
    "method": "支付方式 (1).png",
    "add_card": "添加卡片.png",
    "currency": "货币审查.png",
    "verify": "验证.png",
    "processing": "正在付款中.png",
    "failed": "支付失败.png",
    "complete": "预定完成.png",
}


def set_font(run, size=16, bold=False, color=None):
    run.font.name = FONT_KR
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["deep_text"]


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["bg"]
    bg.line.fill.background()


def add_title(slide, text, left=0.8, top=0.6, width=5.6, height=0.8, size=31):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=True, color=COLORS["deep_text"])
    return tb


def add_body_text(slide, text, left, top, width, height, size=14, color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size=size, color=color or COLORS["soft_text"])
    return tb


def add_rounded_card(slide, left, top, width, height, fill=COLORS["card"], line=COLORS["line"], radius_shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    card = slide.shapes.add_shape(radius_shape, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = line
    card.line.width = Pt(1)
    return card


def add_phone_screenshot(slide, image_path, left, top, width, height, label=None):
    box = add_rounded_card(slide, left, top, width, height, fill=COLORS["white"], line=COLORS["line"])
    inset = 0.06
    ix, iy, iw, ih = left + inset, top + inset, width - inset * 2, height - inset * 2
    p = Path(image_path)
    if p.exists():
        slide.shapes.add_picture(str(p), Inches(ix), Inches(iy), width=Inches(iw), height=Inches(ih))
    else:
        ph = slide.shapes.add_textbox(Inches(ix), Inches(iy), Inches(iw), Inches(ih))
        tf = ph.text_frame
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        run = p0.add_run()
        run.text = f"Placeholder\n{image_path}"
        set_font(run, size=11, color=COLORS["soft_text"])
    if label:
        add_body_text(slide, label, left, top + height + 0.05, width, 0.3, size=10, color=COLORS["soft_text"], align=PP_ALIGN.CENTER)
    return box


def add_footer(slide, page_no):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(7.8), Inches(10.0), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.fill.background()

    add_body_text(slide, "Roami Payment Flow Design", 0.85, 7.86, 3.0, 0.22, size=9, color=COLORS["soft_text"])
    add_body_text(slide, str(page_no), 10.7, 7.84, 0.5, 0.25, size=9, color=COLORS["soft_text"], align=PP_ALIGN.RIGHT)


def add_chip(slide, text, left, top, width=1.55):
    chip = add_rounded_card(slide, left, top, width, 0.45, fill=COLORS["pink"], line=COLORS["line"])
    tf = chip.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_font(run, size=11, bold=True, color=COLORS["deep_text"])


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank); add_bg(s)
    add_title(s, "Roami 결제 경험 고도화 디자인", size=34)
    add_body_text(s, "외국인 관람객을 위한 전시 예약 및 결제 플로우 개선", 0.8, 1.6, 5.2, 0.6, size=16)
    add_body_text(s, "Design Facilitation & Leadership / Personal Project", 0.8, 7.3, 4.8, 0.25, size=9)
    add_phone_screenshot(s, IMAGE_FILES["method"], 6.3, 1.1, 1.45, 4.9)
    add_phone_screenshot(s, IMAGE_FILES["currency"], 7.95, 1.45, 1.45, 4.9)
    add_phone_screenshot(s, IMAGE_FILES["complete"], 9.6, 1.8, 1.45, 4.9)
    add_footer(s, 1)

    #2
    s = prs.slides.add_slide(blank); add_bg(s)
    add_title(s, "외국인 사용자의 결제 불편 문제", size=30)
    add_body_text(s, "외국인 사용자는 한국 전시 예약 과정에서 결제 수단, 금액 이해, 정보 확인, 오류 대응 측면에서 불편을 겪을 수 있다.",0.8,1.45,5.8,0.8,size=13)
    items=["결제 방식의 낯섦","환율 이해의 어려움","예매 정보 확인 부족","결제 실패 후 대처 경로 부족"]
    y=2.35
    for it in items:
        c=add_rounded_card(s,0.8,y,5.2,0.9)
        tf=c.text_frame; tf.clear(); p=tf.paragraphs[0]
        run=p.add_run(); run.text=it; set_font(run,size=16,bold=True)
        y+=1.05
    add_phone_screenshot(s, IMAGE_FILES["failed"], 6.6, 1.5, 4.2, 5.8)
    add_footer(s,2)

    #3
    s=prs.slides.add_slide(blank); add_bg(s)
    add_title(s,"결제 방법 기능의 세부 디자인 방향",size=30)
    add_body_text(s,"기존 Roami의 핵심 기능 중 ‘결제 방법’을 중심으로, 외국인 사용자가 보다 쉽고 안전하게 전시 티켓을 예약할 수 있는 결제 경험을 구체화하였다.",0.8,1.55,5.6,1.3,size=13)
    labels=["언어 선택","실시간 번역","전시 정보","결제 방법"]
    x=6.2
    for i,l in enumerate(labels):
        fill = COLORS["primary"] if i==3 else COLORS["card"]
        card=add_rounded_card(s,x,2.5,1.2,2.2,fill=fill)
        tf=card.text_frame; tf.clear(); tf.vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        run=p.add_run(); run.text=l; set_font(run,size=13,bold=True,color=COLORS["white"] if i==3 else COLORS["deep_text"])
        if i<3:
            add_body_text(s,"→",x+1.22,3.45,0.3,0.3,size=14,color=COLORS["secondary"],align=PP_ALIGN.CENTER)
        x+=1.35
    add_footer(s,3)

    #4
    s=prs.slides.add_slide(blank); add_bg(s)
    add_title(s,"결제 경험 개선을 위한 3가지 목표",size=30)
    goals=[("✓","명확성 (Clarity)","구매 정보와 금액을 쉽게 확인할 수 있도록 한다."),("🔒","신뢰성 (Trust)","인증과 진행 상태 안내를 통해 결제 과정의 안정감을 제공한다."),("↻","회복 가능성 (Recovery)","결제 실패 시 다시 시도하거나 다른 방법을 선택할 수 있도록 한다.")]
    x=0.9
    for icon,t,b in goals:
        c=add_rounded_card(s,x,2.0,3.55,4.2)
        add_body_text(s,icon,x+0.2,2.3,0.6,0.5,size=22,color=COLORS["primary"])
        add_body_text(s,t,x+0.7,2.35,2.7,0.5,size=16,color=COLORS["deep_text"])
        add_body_text(s,b,x+0.25,3.0,3.0,1.8,size=12)
        x+=3.8
    add_footer(s,4)

    #5 flow
    s=prs.slides.add_slide(blank); add_bg(s)
    add_title(s,"전시 티켓 결제를 위한 사용자 플로우",size=30)
    steps=["Select Tickets","Select Date & Time","Order Review","Payment Method","Add New Card","Currency Review","Verification","Processing","Reservation Complete"]
    x=0.55
    for i,st in enumerate(steps):
        w=1.18
        c=add_rounded_card(s,x,3.0,w,0.75,fill=COLORS["card"] if i<8 else COLORS["secondary"])
        tf=c.text_frame; tf.clear(); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        run=p.add_run(); run.text=st; set_font(run,size=8,bold=True)
        if i<8:
            add_body_text(s,"→",x+w+0.02,3.2,0.25,0.3,size=11,color=COLORS["primary"],align=PP_ALIGN.CENTER)
        x+=1.28
    add_body_text(s,"Processing",8.9,4.55,1.0,0.2,size=8,color=COLORS["soft_text"],align=PP_ALIGN.CENTER)
    add_body_text(s,"↓",9.35,4.7,0.2,0.2,size=12,color=COLORS["pink"],align=PP_ALIGN.CENTER)
    fail=add_rounded_card(s,8.25,5.0,1.45,0.72,fill=COLORS["pink"])
    tf=fail.text_frame; tf.clear(); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    run=p.add_run(); run.text="Payment Failed"; set_font(run,size=9,bold=True)
    rec=add_rounded_card(s,9.85,5.0,1.6,0.72,fill=COLORS["pink"])
    tf=rec.text_frame; tf.clear(); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    run=p.add_run(); run.text="Try Again /\nChoose Another Method"; set_font(run,size=8,bold=True)
    add_body_text(s,"결제 기능을 단일 화면이 아닌 전체 사용자 여정으로 확장하여 설계하였다.",0.8,6.5,7.0,0.5,size=12)
    add_footer(s,5)

    #6
    s=prs.slides.add_slide(blank); add_bg(s)
    add_title(s,"티켓과 방문 시간 선택 단계",size=30)
    add_body_text(s,"사용자는 결제 전에 티켓 종류, 수량, 방문 날짜와 시간을 단계적으로 선택할 수 있다. 이를 통해 잘못된 예매를 줄이고 구매 전 정보를 명확하게 인지할 수 있다.",0.8,1.5,5.0,1.6,size=13)
    add_phone_screenshot(s,IMAGE_FILES["ticket"],6.0,1.5,2.25,5.3,label="Ticket Selection")
    add_phone_screenshot(s,IMAGE_FILES["date"],8.55,1.5,2.25,5.3,label="Time Slot Selection")
    add_footer(s,6)

    #7
    s=prs.slides.add_slide(blank); add_bg(s)
    add_title(s,"결제 전 주문 정보 확인 단계",size=30)
    add_phone_screenshot(s,IMAGE_FILES["order"],0.9,1.5,4.9,5.6)
    add_body_text(s,"결제 직전 전시명, 방문 날짜, 시간, 티켓 수량, 연락처와 최종 금액을 한 번 더 확인하도록 구성하였다. 이는 외국인 사용자의 불안감을 줄이고 결제 실수를 예방하기 위한 단계이다.",6.1,1.7,4.7,1.6,size=13)
    for i,t in enumerate(["Exhibition Info","Ticket Details","Price Summary"]):
        c=add_rounded_card(s,6.2,3.5+i*1.15,4.4,0.9)
        tf=c.text_frame; tf.clear(); p=tf.paragraphs[0]
        run=p.add_run(); run.text=t; set_font(run,size=14,bold=True,color=COLORS["primary"])
    add_footer(s,7)

    #8
    s=prs.slides.add_slide(blank); add_bg(s)
    add_title(s,"다양한 결제 수단을 고려한 Payment Method 설계",size=28)
    add_body_text(s,"국제 카드, 모바일 월렛, 아시아 지역 결제 수단을 분류하여 사용자가 익숙한 결제 방법을 쉽게 선택할 수 있도록 하였다. 또한 카드 등록 화면을 통해 새로운 결제 수단을 추가할 수 있게 구성하였다.",0.8,1.5,10.0,1.0,size=12)
    add_chip(s,"International Cards",0.9,2.7,2.2)
    add_chip(s,"Mobile Wallets",3.3,2.7,2.0)
    add_chip(s,"Local & Asian Payment",5.5,2.7,2.5)
    add_phone_screenshot(s,IMAGE_FILES["method"],2.0,3.35,3.4,3.6)
    add_phone_screenshot(s,IMAGE_FILES["add_card"],6.0,3.35,3.4,3.6)
    add_footer(s,8)

    #9
    s=prs.slides.add_slide(blank); add_bg(s)
    add_title(s,"환율 안내와 결제 인증을 통한 신뢰감 형성",size=29)
    add_body_text(s,"외국인 사용자는 원화 금액을 직관적으로 이해하기 어려울 수 있기 때문에, 예상 환율 정보를 제공하여 결제 금액에 대한 이해를 돕는다. 또한 인증 단계와 결제 진행 상태를 시각적으로 제공하여 결제 과정의 신뢰감을 높인다.",0.8,1.4,10.0,1.0,size=12)
    add_phone_screenshot(s,IMAGE_FILES["currency"],1.0,2.5,3.1,4.2,label="Currency Understanding")
    add_phone_screenshot(s,IMAGE_FILES["verify"],4.3,2.5,3.1,4.2,label="Security Verification")
    add_phone_screenshot(s,IMAGE_FILES["processing"],7.6,2.5,3.1,4.2,label="Payment Feedback")
    add_footer(s,9)

    #10
    s=prs.slides.add_slide(blank); add_bg(s)
    add_title(s,"성공과 실패를 모두 고려한 결제 경험 완성",size=29)
    add_body_text(s,"결제 완료 화면뿐만 아니라 결제 실패 화면도 함께 설계하여 사용자가 문제 상황에서도 다음 행동을 쉽게 선택할 수 있도록 하였다. 이를 통해 Roami의 결제 경험은 단순한 구매 기능을 넘어, 외국인 사용자를 위한 안정적이고 포용적인 예약 서비스로 확장되었다.",0.8,1.45,10.0,1.15,size=12)
    add_phone_screenshot(s,IMAGE_FILES["failed"],1.3,2.7,4.2,3.8)
    add_phone_screenshot(s,IMAGE_FILES["complete"],6.2,2.7,4.2,3.8)
    add_chip(s,"명확한 정보",2.2,6.9,2.1)
    add_chip(s,"안전한 결제",4.8,6.9,2.1)
    add_chip(s,"포용적 경험",7.4,6.9,2.1)
    add_footer(s,10)

    out = Path("roami_payment_flow_a4_landscape.pptx")
    prs.save(out)
    print(f"Saved: {out.resolve()}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
